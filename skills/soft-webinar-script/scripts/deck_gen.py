#!/usr/bin/env python3
"""Gera .pptx no ID visual Soft (preto/Bebas/Inter/verde) a partir de um JSON de slides.
Uso: python3 deck_gen.py slides.json saida.pptx
Schema de cada slide:
{ "tipo": "capa|secao|frase|conteudo|prova|mao|oferta|investimento|fechamento",
  "rotulo": "kicker em caps", "titulo": "...", "sub": "...",
  "itens": [ {"titulo":"...", "texto":"..."} ],   # conteudo/oferta/mao
  "numero": "R$ 1.234.567", "fonte": "...",        # prova
  "corpo": "frase grande",                          # frase/secao
  "nota": "speaker notes" }
"""
import json, sys, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG    = RGBColor(0x0A,0x0A,0x0A)
INK   = RGBColor(0xFF,0xFF,0xFF)
DIM   = RGBColor(0xB8,0xB8,0xB8)
FAINT = RGBColor(0x5A,0x5A,0x5A)
GREEN = RGBColor(0x4A,0xDE,0x80)
RED   = RGBColor(0xC0,0x39,0x2B)
LINE  = RGBColor(0x2F,0x2F,0x2F)
TITLEFONT = "Bebas Neue"
BODYFONT  = "Inter"
BRAND = os.environ.get("DECK_BRAND", "IMERSÃO SOFT BUSINESS")

W, H = Inches(13.333), Inches(7.5)
PAD  = Inches(0.7)

def _para(tf, text, size, color, font=BODYFONT, bold=False, align=PP_ALIGN.LEFT, space_after=6, caps=False, line=1.1):
    p = tf.add_paragraph() if tf.paragraphs[0].runs or tf.paragraphs[0].text else tf.paragraphs[0]
    p.alignment = align; p.line_spacing = line; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = (text or "").upper() if caps else (text or "")
    f = r.font; f.size = Pt(size); f.bold = bold; f.name = font; f.color.rgb = color
    return p

def _box(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf

def _bg(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BG

def _kicker(tf, rotulo):
    if rotulo: _para(tf, "—  " + rotulo, 13, DIM, BODYFONT, bold=True, caps=True, space_after=14)

def _footer(slide, rotulo, n):
    tf = _box(slide, PAD, H-Inches(0.5), W-2*PAD, Inches(0.35))
    p = _para(tf, ((BRAND + "   ·   ") + (rotulo or "")).strip(" ·"), 9, FAINT, BODYFONT, caps=True)
    if n is not None:
        tf2 = _box(slide, W-Inches(1.2), H-Inches(0.5), Inches(0.6), Inches(0.35))
        _para(tf2, str(n), 9, FAINT, BODYFONT, align=PP_ALIGN.RIGHT)

def render(slides, out):
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    blank = prs.slide_layouts[6]
    for i, s in enumerate(slides):
        sl = prs.slides.add_slide(blank); _bg(sl)
        tipo = s.get("tipo","conteudo"); rot = s.get("rotulo",""); tit = s.get("titulo","")
        if tipo in ("frase","secao"):
            tf = _box(sl, PAD, Inches(0), W-2*PAD, H, anchor=MSO_ANCHOR.MIDDLE)
            if rot: _para(tf, rot, 13, GREEN, BODYFONT, bold=True, caps=True, space_after=18, align=PP_ALIGN.CENTER)
            _para(tf, s.get("corpo", tit), 46 if tipo=="frase" else 40, INK, TITLEFONT, caps=True, align=PP_ALIGN.CENTER, line=1.0, space_after=0)
        elif tipo == "capa":
            tf = _box(sl, PAD, Inches(0), W-2*PAD, H, anchor=MSO_ANCHOR.MIDDLE)
            _para(tf, rot or "IMERSÃO", 14, GREEN, BODYFONT, bold=True, caps=True, space_after=10, align=PP_ALIGN.CENTER)
            _para(tf, tit, 64, INK, TITLEFONT, caps=True, align=PP_ALIGN.CENTER, line=0.95, space_after=18)
            if s.get("sub"): _para(tf, s["sub"], 16, DIM, BODYFONT, align=PP_ALIGN.CENTER, line=1.3)
        elif tipo == "prova":
            tf = _box(sl, PAD, Inches(0), W-2*PAD, H, anchor=MSO_ANCHOR.MIDDLE)
            _para(tf, s.get("numero",""), 96, GREEN, TITLEFONT, align=PP_ALIGN.CENTER, space_after=10, line=0.9)
            if tit: _para(tf, tit, 20, INK, BODYFONT, align=PP_ALIGN.CENTER, space_after=6)
            if s.get("fonte"): _para(tf, s["fonte"], 13, DIM, BODYFONT, align=PP_ALIGN.CENTER)
        else:  # conteudo / mao / oferta / investimento / fechamento
            tf = _box(sl, PAD, PAD, W-2*PAD, Inches(2.0))
            _kicker(tf, rot)
            _para(tf, tit, 40, INK, TITLEFONT, caps=True, line=0.98, space_after=4)
            if s.get("sub"): _para(tf, s["sub"], 16, DIM, BODYFONT, space_after=0)
            body = _box(sl, PAD, PAD+Inches(2.05), W-2*PAD, H-PAD-Inches(2.6))
            for it in s.get("itens", []):
                if isinstance(it, dict):
                    _para(body, it.get("titulo",""), 17, GREEN if tipo in ("oferta","investimento") else INK, BODYFONT, bold=True, space_after=2)
                    if it.get("texto"): _para(body, it["texto"], 14, DIM, BODYFONT, space_after=12, line=1.25)
                else:
                    _para(body, str(it), 16, DIM, BODYFONT, space_after=10, line=1.3)
            if s.get("corpo"): _para(body, s["corpo"], 16, DIM, BODYFONT, line=1.4)
        _footer(sl, rot, i+1 if tipo not in ("capa",) else None)
        if s.get("nota"):
            sl.notes_slide.notes_text_frame.text = s["nota"]
    prs.save(out); print(f"OK {len(slides)} slides -> {out}")

if __name__ == "__main__":
    src = sys.argv[1] if len(sys.argv) > 1 else "slides.json"
    out = sys.argv[2] if len(sys.argv) > 2 else "deck.pptx"
    with open(src, encoding="utf-8") as fh: render(json.load(fh), out)
